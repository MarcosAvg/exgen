from pptx import Presentation
from pptx.util import Inches, Pt
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
table = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(5), Inches(2)).table
table.columns[0].width = Pt(50)
table.columns[1].width = Pt(310)
prs.save('test_table.pptx')
